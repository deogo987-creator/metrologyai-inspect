#!/usr/bin/env python3
"""
Edit the SIH 2026 template directly.
Preserves ALL shapes, images, positions, colors, and formatting.
Only replaces content text in the existing text boxes.
"""

from pptx import Presentation
from pptx.util import Pt
from copy import deepcopy

INPUT = "SIH2026-IDEA-Presentation-Format.pptx"
OUTPUT = "MetrologyAI_SIH_2026_Final.pptx"

prs = Presentation(INPUT)

# ──────────────────────────────────────────────────────────
# SLIDE 1: Title Page — edit shape[5] (PS ID/Title bullets)
# ──────────────────────────────────────────────────────────
slide1 = prs.slides[0]
shape5 = slide1.shapes[5]  # Text box with PS details

# Clear existing paragraphs and replace with our content
# We need to preserve the first paragraph and clear the rest
tf = shape5.text_frame

# Get the first run's font properties to preserve formatting
first_para = tf.paragraphs[0]
first_run = first_para.runs[0] if first_para.runs else None
font_size = Pt(24)
font_bold = True

# Set new text for each paragraph
new_lines = [
    "Problem Statement ID – PS26034",
    "Problem Statement Title – AI-Assisted Legal Metrology Label Compliance Inspection System",
    "Theme – Software",
    "PS Category – Software",
    "Team ID – [TEAM ID]",
    "Team Name – [TEAM NAME] (Registered on portal)"
]

# Clear all paragraphs first
for i in range(len(tf.paragraphs) - 1, 0, -1):
    p = tf.paragraphs[i]
    p.clear()

# Now set text on each paragraph
for i, line in enumerate(new_lines):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()

    p.clear()
    run = p.add_run()
    run.text = line
    run.font.size = font_size
    run.font.bold = font_bold
    # Preserve scheme color (no explicit rgb = inherits from theme)

# ──────────────────────────────────────────────────────────
# SLIDE 2: IDEA TITLE → METROLOGYAI
# Edit shape[1] title + shape[2] content
# ──────────────────────────────────────────────────────────
slide2 = prs.slides[1]

# Edit title (shape[1]) — replace "IDEA TITLE" with "METROLOGYAI"
title_shape = slide2.shapes[1]
tf = title_shape.text_frame
for p in tf.paragraphs:
    for run in p.runs:
        if "IDEA TITLE" in run.text:
            run.text = "METROLOGYAI"

# Edit content (shape[2]) — replace placeholder bullets
content_shape = slide2.shapes[2]
tf = content_shape.text_frame

# Preserve font from first run
first_run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else None

# Clear all paragraphs
for i in range(len(tf.paragraphs) - 1, 0, -1):
    p = tf.paragraphs[i]
    p.clear()

# New content for Slide 2
slide2_content = [
    ("Proposed Solution — AI-Assisted Legal Metrology Compliance Inspection System", True, Pt(32)),
    ("MetrologyAI is a software system that assists inspectors in checking packaged commodity labels under the Legal Metrology (Packaged Commodities) Rules, 2011.", False, Pt(28)),
    ("The system uses product images, OCR, computer vision and a configurable compliance rule engine to automate label analysis — the inspector remains responsible for the final verification and decision.", False, Pt(28)),
    ("AI assists with: image analysis, text extraction, declaration identification, compliance validation, evidence generation, risk prioritization, cross-source comparison, inspection recommendations and digital reporting.", False, Pt(28)),
    ("Unique capabilities: Multi-View Intelligence (front, back, side, top, bottom), Missing ≠ Unreadable distinction, Adaptive Re-Capture, Evidence Chain linking, Cross-Source Verification and Human-in-the-Loop revalidation.", False, Pt(28)),
]

for i, (text, bold, size) in enumerate(slide2_content):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()

    p.clear()
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold

# ──────────────────────────────────────────────────────────
# SLIDE 3: TECHNICAL APPROACH — edit shape[2] content
# ──────────────────────────────────────────────────────────
slide3 = prs.slides[2]
content_shape = slide3.shapes[2]
tf = content_shape.text_frame

# Clear all paragraphs
for i in range(len(tf.paragraphs) - 1, 0, -1):
    p = tf.paragraphs[i]
    p.clear()

slide3_content = [
    ("Technologies used: React, TypeScript, Vite, Tailwind CSS (Frontend) • Convex Serverless Backend & Database • Gemini AI Vision + OCR (AI/ML) • Recharts (Analytics)", False, Pt(28)),
    ("Methodology: Multi-view product image capture → Image quality analysis (blur, glare, visibility) → OCR + AI vision text extraction → Structured declaration extraction (MRP, manufacturer, net quantity, consumer care, country of origin) → Compliance rule validation → Evidence engine linking findings to source → Risk prioritization → Inspector review → Digital report generation", False, Pt(28)),
    ("Implementation flow: Dashboard → New Inspection → Upload product images (multi-view) → AI analysis → Compliance results with evidence chains → Inspector corrections → Smart revalidation → Export report", False, Pt(28)),
    ("Architecture: Modular design separating AI extraction from regulatory rules, enabling independent updates to vision models and compliance rule sets without system-wide changes", False, Pt(28)),
    ("Deployment: Web-based application supporting desktop, tablet and mobile workflows for field inspection scenarios", False, Pt(28)),
]

for i, (text, bold, size) in enumerate(slide3_content):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.clear()
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold

# ──────────────────────────────────────────────────────────
# SLIDE 4: FEASIBILITY AND VIABILITY — edit shape[2] content
# ──────────────────────────────────────────────────────────
slide4 = prs.slides[3]
content_shape = slide4.shapes[2]
tf = content_shape.text_frame

for i in range(len(tf.paragraphs) - 1, 0, -1):
    p = tf.paragraphs[i]
    p.clear()

slide4_content = [
    ("Technology Feasibility: OCR, computer vision and AI vision models (Gemini) are technically mature and applicable to label inspection tasks with documented accuracy in scene-text recognition.", False, Pt(28)),
    ("Data Feasibility: Product-label images and structured declaration fields (MRP, manufacturer, net quantity) can be systematically collected for development, testing and continuous improvement.", False, Pt(28)),
    ("Architecture Feasibility: AI extraction and regulatory compliance rules are separated into modular components, allowing independent updates to vision models and rule definitions.", False, Pt(28)),
    ("Deployment Feasibility: A web-based system (React + Convex) supports desktop, tablet and mobile workflows required for field inspection scenarios.", False, Pt(28)),
    ("Scalability: New product categories, declaration types and rule sets can be added through configurable definitions without code changes to the core system.", False, Pt(28)),
    ("Challenges and Mitigations: AI misreading text → Confidence scoring + evidence + inspector correction • Blurry images → Image-quality detection + adaptive recapture • Declarations on other sides → Multi-view package analysis • Online information differs → Cross-source comparison • Regulations change → Version-controlled rule engine • False AI findings → Human-in-the-loop verification", False, Pt(28)),
]

for i, (text, bold, size) in enumerate(slide4_content):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.clear()
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold

# ──────────────────────────────────────────────────────────
# SLIDE 5: IMPACT AND BENEFITS — edit shape[2] content
# ──────────────────────────────────────────────────────────
slide5 = prs.slides[4]
content_shape = slide5.shapes[2]
tf = content_shape.text_frame

for i in range(len(tf.paragraphs) - 1, 0, -1):
    p = tf.paragraphs[i]
    p.clear()

slide5_content = [
    ("Potential Impact on Target Audience: Faster inspection through automated label-reading and initial validation; consistent checking via systematic rule-based validation; better evidence through image-linked findings; reduced false positives by distinguishing missing from unreadable declarations; improved decision support with recommended next actions for inspectors.", False, Pt(28)),
    ("Inspector and Operational Benefits: AI-assisted extraction reduces repetitive manual work; evidence-based findings provide traceable inspection records; risk prioritization focuses inspector attention on high-priority violations; digital reports replace manual documentation; human-in-the-loop ensures inspector retains final authority.", False, Pt(28)),
    ("Economic and Strategic Benefits: Reduced manual effort in repetitive label verification; scalable workflow supporting larger inspection volumes; structured digital documentation for audit and compliance records; data-driven monitoring enabling analysis of recurring findings across inspections; future expansion to additional product categories and rule sets.", False, Pt(28)),
    ("Social Benefits: Improved consumer protection through more consistent compliance checking; enhanced transparency in inspection processes; support for Legal Metrology enforcement at scale; digital transformation of government inspection workflows.", False, Pt(28)),
]

for i, (text, bold, size) in enumerate(slide5_content):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.clear()
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold

# ──────────────────────────────────────────────────────────
# SLIDE 6: RESEARCH AND REFERENCES — edit shape[2] content
# ──────────────────────────────────────────────────────────
slide6 = prs.slides[5]
content_shape = slide6.shapes[2]
tf = content_shape.text_frame

for i in range(len(tf.paragraphs) - 1, 0, -1):
    p = tf.paragraphs[i]
    p.clear()

slide6_content = [
    ("Legal and Regulatory References: Legal Metrology Act, 2009 • Legal Metrology (Packaged Commodities) Rules, 2011 • Department of Consumer Affairs — Legal Metrology Division • Relevant official notifications and amendments under the Legal Metrology framework", False, Pt(28)),
    ("Research Areas: Computer Vision and OCR for scene-text recognition in natural environments • Multimodal AI for combined image and text understanding • Explainable AI for confidence scoring and evidence-based decision support • Human-AI Interaction for verified inspection workflows", False, Pt(28)),
    ("Technology References: React (UI library) • TypeScript (type-safe development) • Convex (serverless backend and database) • Tailwind CSS (styling) • Gemini AI Vision (image analysis and OCR) • Vite (build tool) • Recharts (analytics visualization)", False, Pt(28)),
    ("References: [ADD VERIFIED REFERENCE — Legal Metrology Act, 2009] • [ADD VERIFIED REFERENCE — Legal Metrology (Packaged Commodities) Rules, 2011] • [ADD VERIFIED REFERENCE — Department of Consumer Affairs, Government of India] • [ADD VERIFIED REFERENCE — Relevant court judgments or regulatory orders]", False, Pt(28)),
]

for i, (text, bold, size) in enumerate(slide6_content):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.clear()
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold

# ──────────────────────────────────────────────────────────
# SAVE — keep Slide 7 (instructions) unchanged
# ──────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"✅ Edited template saved to: {OUTPUT}")
print(f"   Slides: {len(prs.slides)} (6 content + 1 instructions)")
print(f"   All shapes, images, positions, colors preserved")
