#!/usr/bin/env python3
"""
MetrologyAI - SIH 2026 Presentation Generator
Generates a 6-slide PPT matching the SIH 2026 template theme.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ── Color Palette (matching SIH 2026 template) ──
NAVY = RGBColor(0x1A, 0x3C, 0x6E)        # Dark navy for headers
BLUE = RGBColor(0x00, 0x56, 0xA0)        # SIH blue
LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)  # Light blue background
ORANGE = RGBColor(0xE8, 0x6C, 0x00)      # SIH orange
GREEN = RGBColor(0x1B, 0x8A, 0x4E)       # SIH green
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
DARK_GRAY = RGBColor(0x55, 0x55, 0x55)
MEDIUM_GRAY = RGBColor(0x99, 0x99, 0x99)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)    # Slide background tint
FOOTER_BLUE = RGBColor(0x00, 0x56, 0xA0) # Footer bar
CARD_BG = RGBColor(0xF0, 0xF4, 0xF8)     # Card background
ACCENT_BLUE = RGBColor(0x33, 0x7A, 0xB7) # Accent for shapes
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
ACCENT_TEAL = RGBColor(0x16, 0xA0, 0x85)

# ── Slide dimensions (widescreen 16:9) ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]


def add_footer(slide, page_num):
    """Add blue footer bar matching SIH template."""
    footer = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.05),
        SLIDE_W, Inches(0.45)
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = FOOTER_BLUE
    footer.line.fill.background()

    tf = footer.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "@SIH Idea submission - MetrologyAI"
    run.font.size = Pt(10)
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"

    # Page number
    num_box = slide.shapes.add_textbox(
        Inches(12.2), Inches(7.08),
        Inches(1), Inches(0.4)
    )
    tf2 = num_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run2 = p2.add_run()
    run2.text = str(page_num)
    run2.font.size = Pt(12)
    run2.font.color.rgb = WHITE
    run2.font.bold = True
    run2.font.name = "Calibri"


def add_team_name_oval(slide):
    """Add 'Your Team Name' oval in top-left corner."""
    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(0.3), Inches(0.15),
        Inches(1.3), Inches(0.85)
    )
    oval.fill.background()
    oval.line.color.rgb = NAVY
    oval.line.width = Pt(1.5)

    tf = oval.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Your\nTeam Name"
    run.font.size = Pt(10)
    run.font.color.rgb = NAVY
    run.font.bold = True
    run.font.name = "Calibri"


def add_sih_logo(slide):
    """Add SIH logo text placeholder in top-right."""
    # SIH brain logo text representation
    logo_box = slide.shapes.add_textbox(
        Inches(11.2), Inches(0.05),
        Inches(2), Inches(1.0)
    )
    tf = logo_box.text_frame
    tf.word_wrap = True

    # "SMART INDIA" line
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.RIGHT
    r1 = p1.add_run()
    r1.text = "SMART INDIA"
    r1.font.size = Pt(11)
    r1.font.color.rgb = NAVY
    r1.font.bold = True
    r1.font.name = "Calibri"

    # "HACKATHON 2026" line
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = "HACKATHON 2026"
    r2.font.size = Pt(11)
    r2.font.color.rgb = GREEN
    r2.font.bold = True
    r2.font.name = "Calibri"

    # SIH badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(11.8), Inches(0.7),
        Inches(0.6), Inches(0.35)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = ORANGE
    badge.line.fill.background()
    btf = badge.text_frame
    btf.paragraphs[0].alignment = PP_ALIGN.CENTER
    brun = btf.paragraphs[0].add_run()
    brun.text = "SIH"
    brun.font.size = Pt(8)
    brun.font.color.rgb = WHITE
    brun.font.bold = True
    brun.font.name = "Calibri"


def add_header_text(slide, title_text, subtitle_text=None):
    """Add title text centered below header area."""
    # Main title
    title_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(0.2),
        Inches(10), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(28)
    run.font.color.rgb = NAVY
    run.font.bold = True
    run.font.name = "Calibri"

    if subtitle_text:
        sub_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(0.85),
            Inches(10), Inches(0.4)
        )
        stf = sub_box.text_frame
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        srun = sp.add_run()
        srun.text = subtitle_text
        srun.font.size = Pt(14)
        srun.font.color.rgb = DARK_GRAY
        srun.font.italic = True
        srun.font.name = "Calibri"


def add_rounded_card(slide, left, top, width, height, fill_color=CARD_BG, border_color=None):
    """Add a rounded rectangle card."""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
    else:
        card.line.fill.background()
    return card


def add_text_box(slide, left, top, width, height, text, font_size=10,
                 color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, italic=False,
                 font_name="Calibri"):
    """Add a text box with formatted text."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    return tb


def add_multiline_text(slide, left, top, width, height, lines, font_size=10,
                       color=BLACK, bold_first=False, line_spacing=1.2, bullet=False):
    """Add multiple lines of text."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.space_after = Pt(2)
        p.space_before = Pt(1)

        if bullet:
            prefix = "• "
        else:
            prefix = ""

        run = p.add_run()
        run.text = prefix + line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
        if bold_first and i == 0:
            run.font.bold = True

    return tb


def add_arrow_right(slide, left, top, width=Inches(0.3), height=Inches(0.2)):
    """Add right-pointing arrow."""
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        left, top, width, height
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = BLUE
    arrow.line.fill.background()
    return arrow


def add_arrow_down(slide, left, top, width=Inches(0.2), height=Inches(0.25)):
    """Add down-pointing arrow."""
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        left, top, width, height
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = BLUE
    arrow.line.fill.background()
    return arrow


def add_connector_line(slide, x1, y1, x2, y2, color=BLUE, width=Pt(2)):
    """Add a simple connector line using a thin rectangle."""
    # Calculate angle and length
    dx = x2 - x1
    dy = y2 - y1
    length = int(math.sqrt(dx*dx + dy*dy))
    angle = math.atan2(dy, dx)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        x1, y1,
        length, Pt(2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    line.rotation = math.degrees(angle)
    return line


# ══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE PAGE — Understanding the Problem + Solution
# ══════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(blank_layout)

# SIH Header
add_sih_logo(slide1)

# "SMART INDIA HACKATHON 2026" header
hdr = slide1.shapes.add_textbox(Inches(1.5), Inches(0.05), Inches(9), Inches(0.45))
htf = hdr.text_frame
hp = htf.paragraphs[0]
hp.alignment = PP_ALIGN.CENTER
hr = hp.add_run()
hr.text = "SMART INDIA HACKATHON 2026"
hr.font.size = Pt(22)
hr.font.color.rgb = NAVY
hr.font.bold = True
hr.font.name = "Calibri"

# Team name oval
add_team_name_oval(slide1)

# ── LEFT SECTION: Understanding the Problem ──
add_text_box(slide1, Inches(0.3), Inches(1.15), Inches(4.2), Inches(0.4),
             "UNDERSTANDING THE PROBLEM", 14, NAVY, True, PP_ALIGN.CENTER)

problems = [
    ("01", "Manual Inspection", "Inspectors manually examine labels, making repetitive checks time-consuming."),
    ("02", "Multi-Surface Declarations", "Required info may appear on different sides of the package."),
    ("03", "Real-World Image Issues", "Blur, glare, distortion and small text reduce readability."),
    ("04", "Uncertain AI Results", "A system may confuse unreadable with genuinely missing declarations."),
    ("05", "Fragmented Verification", "Physical packaging, online listings and history need separate comparison."),
]

y_start = Inches(1.6)
for i, (num, title, desc) in enumerate(problems):
    y = y_start + Inches(i * 0.95)
    # Number circle
    circ = slide1.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.3), y, Inches(0.35), Inches(0.35)
    )
    circ.fill.solid()
    circ.fill.fore_color.rgb = BLUE
    circ.line.fill.background()
    ctf = circ.text_frame
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    cr = ctf.paragraphs[0].add_run()
    cr.text = num
    cr.font.size = Pt(8)
    cr.font.color.rgb = WHITE
    cr.font.bold = True
    cr.font.name = "Calibri"

    add_text_box(slide1, Inches(0.75), y - Inches(0.02), Inches(3.7), Inches(0.3),
                 title, 10, NAVY, True)
    add_text_box(slide1, Inches(0.75), y + Inches(0.22), Inches(3.7), Inches(0.55),
                 desc, 8, DARK_GRAY)

# Challenge quote
add_text_box(slide1, Inches(0.3), Inches(6.4), Inches(4.2), Inches(0.5),
             "The challenge is not simply reading a label — it is producing a reliable, explainable inspection decision.",
             8, ACCENT_BLUE, italic=True, alignment=PP_ALIGN.CENTER)

# ── CENTER SECTION: Our Solution Pipeline ──
add_text_box(slide1, Inches(4.6), Inches(1.15), Inches(4.2), Inches(0.4),
             "OUR SOLUTION — METROLOGYAI", 14, NAVY, True, PP_ALIGN.CENTER)

pipeline_steps = [
    "PRODUCT IMAGE",
    "IMAGE QUALITY CHECK",
    "AI VISION + OCR",
    "DECLARATION EXTRACTION",
    "RULE VALIDATION",
    "EVIDENCE + CONFIDENCE",
    "RISK PRIORITIZATION",
    "INSPECTOR REVIEW",
    "DIGITAL REPORT",
]

pipe_x = Inches(5.2)
pipe_start_y = Inches(1.65)
step_h = Inches(0.38)
gap = Inches(0.12)

for i, step in enumerate(pipeline_steps):
    y = pipe_start_y + i * (step_h + gap)
    # Step box
    box = slide1.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        pipe_x, y, Inches(3.0), step_h
    )
    # Color gradient effect - alternate colors
    colors = [BLUE, ACCENT_TEAL, GREEN, ACCENT_BLUE, ACCENT_ORANGE, ACCENT_GREEN, ORANGE, NAVY, BLUE]
    box.fill.solid()
    box.fill.fore_color.rgb = colors[i]
    box.line.fill.background()
    btf = box.text_frame
    btf.paragraphs[0].alignment = PP_ALIGN.CENTER
    brun = btf.paragraphs[0].add_run()
    brun.text = step
    brun.font.size = Pt(8)
    brun.font.color.rgb = WHITE
    brun.font.bold = True
    brun.font.name = "Calibri"
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Arrow between steps
    if i < len(pipeline_steps) - 1:
        add_arrow_down(slide1, Inches(6.55), y + step_h - Inches(0.02), Inches(0.18), Inches(0.14))

# Solution quote
add_text_box(slide1, Inches(4.6), Inches(6.4), Inches(4.2), Inches(0.5),
             "AI assists the inspection; the inspector makes the final decision.",
             8, ACCENT_BLUE, italic=True, alignment=PP_ALIGN.CENTER)

# ── RIGHT SECTION: Unique Solutions ──
add_text_box(slide1, Inches(8.8), Inches(1.15), Inches(4.2), Inches(0.4),
             "OUR UNIQUE SOLUTIONS", 14, NAVY, True, PP_ALIGN.CENTER)

features = [
    ("Multi-View Intelligence", "Analyzes front, back, side, top and bottom views."),
    ("Missing ≠ Unreadable", "Separates absent declarations from low-confidence detection."),
    ("Adaptive Re-Capture", "Identifies poor image quality and recommends another capture."),
    ("Evidence Chain", "Links every finding to image evidence, OCR, field and rule."),
    ("Cross-Source Verification", "Compares physical packaging with online listings."),
    ("Human-in-the-Loop", "Inspector corrections trigger smart revalidation."),
]

feat_y = Inches(1.6)
for i, (title, desc) in enumerate(features):
    y = feat_y + Inches(i * 0.95)
    # Feature card
    card = add_rounded_card(slide1, Inches(8.8), y, Inches(4.2), Inches(0.85),
                           CARD_BG, ACCENT_BLUE)
    # Feature icon dot
    dot_colors = [BLUE, GREEN, ORANGE, ACCENT_TEAL, ACCENT_RED, NAVY]
    dot = slide1.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(8.95), y + Inches(0.12), Inches(0.22), Inches(0.22)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = dot_colors[i]
    dot.line.fill.background()

    add_text_box(slide1, Inches(9.25), y + Inches(0.05), Inches(3.5), Inches(0.3),
                 title, 9, NAVY, True)
    add_text_box(slide1, Inches(9.25), y + Inches(0.35), Inches(3.5), Inches(0.45),
                 desc, 7.5, DARK_GRAY)

# ── BOTTOM: Prototype ──
proto_box = add_rounded_card(slide1, Inches(0.3), Inches(6.85), Inches(12.7), Inches(0.15),
                              BLUE)

# PS ID and Team info bar
info_bar = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(6.85),
    SLIDE_W, Inches(0.2)
)
info_bar.fill.solid()
info_bar.fill.fore_color.rgb = NAVY
info_bar.line.fill.background()
itf = info_bar.text_frame
itf.paragraphs[0].alignment = PP_ALIGN.CENTER
ir = itf.paragraphs[0].add_run()
ir.text = "PS ID: PS26034  |  Team ID: [TEAM ID]  |  Team Name: [TEAM NAME]  |  Category: Software"
ir.font.size = Pt(9)
ir.font.color.rgb = WHITE
ir.font.name = "Calibri"

add_footer(slide1, 1)


# ══════════════════════════════════════════════════════════════
# SLIDE 2: METHODOLOGY + TECHNICAL APPROACH + TECH STACK
# ══════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(blank_layout)
add_sih_logo(slide2)
add_team_name_oval(slide2)
add_header_text(slide2, "METHODOLOGY & TECHNICAL APPROACH")

# ── TOP: Methodology Pipeline ──
add_text_box(slide2, Inches(0.3), Inches(1.15), Inches(12.7), Inches(0.35),
             "METHODOLOGY — End-to-End Inspection Pipeline", 13, NAVY, True, PP_ALIGN.CENTER)

method_steps = [
    ("01", "CAPTURE", "Multi-view\nproduct images", BLUE),
    ("02", "QUALITY\nCHECK", "Blur + glare +\nvisibility", ACCENT_TEAL),
    ("03", "EXTRACT", "OCR + AI\nvision", GREEN),
    ("04", "STRUCTURE", "Declaration\nextraction", ACCENT_BLUE),
    ("05", "VALIDATE", "Compliance\nrule engine", ORANGE),
    ("06", "REVIEW", "Inspector\nverification", NAVY),
    ("07", "REPORT", "Evidence-backed\nreport", ACCENT_GREEN),
]

step_w = Inches(1.55)
step_h_m = Inches(0.9)
start_x = Inches(0.4)
method_y = Inches(1.55)

for i, (num, title, desc, color) in enumerate(method_steps):
    x = start_x + i * (step_w + Inches(0.2))

    card = add_rounded_card(slide2, x, method_y, step_w, step_h_m, color)

    # Number badge
    badge = slide2.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Inches(0.05), method_y + Inches(0.05),
        Inches(0.3), Inches(0.3)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = WHITE
    badge.line.fill.background()
    btf = badge.text_frame
    btf.paragraphs[0].alignment = PP_ALIGN.CENTER
    brun = btf.paragraphs[0].add_run()
    brun.text = num
    brun.font.size = Pt(9)
    brun.font.color.rgb = color
    brun.font.bold = True
    brun.font.name = "Calibri"

    add_text_box(slide2, x + Inches(0.02), method_y + Inches(0.08), step_w - Inches(0.04), Inches(0.25),
                 title, 8, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide2, x + Inches(0.02), method_y + Inches(0.4), step_w - Inches(0.04), Inches(0.45),
                 desc, 7, RGBColor(0xE0, 0xE0, 0xE0), alignment=PP_ALIGN.CENTER)

    # Arrow between steps
    if i < len(method_steps) - 1:
        add_arrow_right(slide2, x + step_w - Inches(0.05), method_y + Inches(0.35),
                       Inches(0.22), Inches(0.18))

# ── LEFT: Technical Approach ──
add_text_box(slide2, Inches(0.3), Inches(2.65), Inches(5.5), Inches(0.35),
             "TECHNICAL APPROACH", 13, NAVY, True, PP_ALIGN.LEFT)

tech_blocks = [
    ("Image Acquisition", "Capture product package images and optional e-commerce listing screenshots.", BLUE),
    ("Image Processing", "Analyze blur, glare, perspective distortion and image quality.", ACCENT_TEAL),
    ("OCR + Vision", "Extract text and identify relevant label regions.", GREEN),
    ("Structured Extraction", "Convert raw text into fields: Product name, MRP, Net quantity, Manufacturer, Consumer care, Country of origin.", ACCENT_BLUE),
    ("Compliance Engine", "Validate extracted information against configurable rules.", ORANGE),
    ("Evidence Engine", "Connect findings to their source image and detected region.", NAVY),
    ("Decision Support", "Generate confidence, priority and recommended next action.", ACCENT_GREEN),
]

tech_y = Inches(3.05)
for i, (title, desc, color) in enumerate(tech_blocks):
    y = tech_y + Inches(i * 0.55)

    # Color bar
    bar = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.3), y, Inches(0.08), Inches(0.45)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text_box(slide2, Inches(0.5), y, Inches(2.0), Inches(0.22), title, 8.5, NAVY, True)
    add_text_box(slide2, Inches(0.5), y + Inches(0.2), Inches(5.0), Inches(0.3), desc, 7, DARK_GRAY)

# ── RIGHT: UI Flow ──
add_text_box(slide2, Inches(6.0), Inches(2.65), Inches(7.0), Inches(0.35),
             "USER INTERFACE — IMPLEMENTATION FLOW", 13, NAVY, True, PP_ALIGN.LEFT)

ui_screens = [
    ("SCREEN 01", "Dashboard"),
    ("SCREEN 02", "New Inspection"),
    ("SCREEN 03", "Image Upload"),
    ("SCREEN 04", "AI Analysis"),
    ("SCREEN 05", "Compliance Result"),
    ("SCREEN 06", "Report"),
]

screen_y = Inches(3.05)
for i, (label, name) in enumerate(ui_screens):
    y = screen_y + Inches(i * 0.55)

    # Screen placeholder card
    card = add_rounded_card(slide2, Inches(6.2), y, Inches(3.0), Inches(0.45),
                           CARD_BG, ACCENT_BLUE)
    add_text_box(slide2, Inches(6.3), y + Inches(0.02), Inches(1.2), Inches(0.2),
                 label, 7, ACCENT_BLUE, True)
    add_text_box(slide2, Inches(7.5), y + Inches(0.02), Inches(1.5), Inches(0.2),
                 name, 9, NAVY, True)
    add_text_box(slide2, Inches(6.3), y + Inches(0.22), Inches(2.8), Inches(0.2),
                 "[INSERT PROTOTYPE SCREENSHOT]", 7, MEDIUM_GRAY, italic=True)

    # Arrow
    if i < len(ui_screens) - 1:
        add_arrow_down(slide2, Inches(7.55), y + Inches(0.42), Inches(0.15), Inches(0.12))

# ── BOTTOM: Tech Stack ──
add_text_box(slide2, Inches(0.3), Inches(6.55), Inches(12.7), Inches(0.3),
             "TECH STACK", 12, NAVY, True, PP_ALIGN.CENTER)

stack_categories = [
    ("FRONTEND", "React • TypeScript • Vite • Tailwind CSS", BLUE),
    ("BACKEND", "Convex (Serverless)", ACCENT_TEAL),
    ("DATABASE", "Convex DB", GREEN),
    ("AI / VISION", "Gemini AI Vision • OCR", ORANGE),
    ("ANALYTICS", "Recharts", NAVY),
    ("DEPLOYMENT", "Freebuff Platform", ACCENT_GREEN),
]

stack_y = Inches(6.85)
stack_w = Inches(2.05)
for i, (cat, techs, color) in enumerate(stack_categories):
    x = Inches(0.3) + i * (stack_w + Inches(0.1))

    cat_box = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, stack_y, stack_w, Inches(0.15)
    )
    cat_box.fill.solid()
    cat_box.fill.fore_color.rgb = color
    cat_box.line.fill.background()
    ctf = cat_box.text_frame
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    cr = ctf.paragraphs[0].add_run()
    cr.text = f"{cat}: {techs}"
    cr.font.size = Pt(6.5)
    cr.font.color.rgb = WHITE
    cr.font.bold = True
    cr.font.name = "Calibri"

add_footer(slide2, 2)


# ══════════════════════════════════════════════════════════════
# SLIDE 3: IMPACT & BENEFITS
# ══════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(blank_layout)
add_sih_logo(slide3)
add_team_name_oval(slide3)
add_header_text(slide3, "IMPACT & BENEFITS")

# ── LEFT: Potential Impact ──
add_text_box(slide3, Inches(0.3), Inches(1.15), Inches(4.0), Inches(0.35),
             "POTENTIAL IMPACT", 13, NAVY, True, PP_ALIGN.LEFT)

impacts = [
    ("Faster Inspection", "Automates repetitive label-reading and initial validation."),
    ("Consistent Checking", "Applies configured validation logic systematically."),
    ("Better Evidence", "Links findings with their supporting image evidence."),
    ("Reduced False Positives", "Distinguishes missing info from unreadable info."),
    ("Improved Decision Support", "Provides inspectors with recommended next actions."),
]

impact_y = Inches(1.6)
for i, (title, desc) in enumerate(impacts):
    y = impact_y + Inches(i * 0.95)

    icon_colors = [BLUE, GREEN, ACCENT_TEAL, ORANGE, NAVY]
    dot = slide3.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.3), y + Inches(0.05), Inches(0.25), Inches(0.25)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = icon_colors[i]
    dot.line.fill.background()

    add_text_box(slide3, Inches(0.65), y, Inches(3.6), Inches(0.25), title, 10, NAVY, True)
    add_text_box(slide3, Inches(0.65), y + Inches(0.3), Inches(3.6), Inches(0.5), desc, 8, DARK_GRAY)

# ── CENTER: Inspector Benefits Hub ──
add_text_box(slide3, Inches(4.5), Inches(1.15), Inches(4.5), Inches(0.35),
             "INSPECTOR & OPERATIONAL BENEFITS", 13, NAVY, True, PP_ALIGN.CENTER)

# Center hub
hub = slide3.shapes.add_shape(
    MSO_SHAPE.OVAL, Inches(6.1), Inches(3.2), Inches(1.5), Inches(1.5)
)
hub.fill.solid()
hub.fill.fore_color.rgb = BLUE
hub.line.fill.background()
htf = hub.text_frame
htf.paragraphs[0].alignment = PP_ALIGN.CENTER
hr = htf.paragraphs[0].add_run()
hr.text = "INSPECTOR"
hr.font.size = Pt(12)
hr.font.color.rgb = WHITE
hr.font.bold = True
hr.font.name = "Calibri"

# Spokes
spokes = [
    ("Faster Capture", Inches(5.2), Inches(2.0)),
    ("AI-Assisted\nExtraction", Inches(7.8), Inches(2.0)),
    ("Evidence-Based\nFindings", Inches(4.6), Inches(3.8)),
    ("Risk\nPrioritization", Inches(8.4), Inches(3.8)),
    ("Digital\nReports", Inches(5.2), Inches(5.2)),
    ("Human\nVerification", Inches(7.8), Inches(5.2)),
]

for label, sx, sy in spokes:
    spoke = add_rounded_card(slide3, sx, sy, Inches(1.3), Inches(0.6), LIGHT_BLUE, ACCENT_BLUE)
    stf = spoke.text_frame
    stf.word_wrap = True
    stf.paragraphs[0].alignment = PP_ALIGN.CENTER
    sr = stf.paragraphs[0].add_run()
    sr.text = label
    sr.font.size = Pt(8)
    sr.font.color.rgb = NAVY
    sr.font.bold = True
    sr.font.name = "Calibri"
    stf.vertical_anchor = MSO_ANCHOR.MIDDLE

# ── RIGHT: Economic & Strategic ──
add_text_box(slide3, Inches(9.5), Inches(1.15), Inches(3.8), Inches(0.35),
             "ECONOMIC & STRATEGIC BENEFITS", 13, NAVY, True, PP_ALIGN.LEFT)

econ_benefits = [
    ("Reduced Manual Effort", "Less repetitive label verification."),
    ("Scalable Workflow", "Supports larger inspection volumes."),
    ("Digital Documentation", "Structured inspection records and reports."),
    ("Data-Driven Monitoring", "Enables analysis of recurring findings."),
    ("Future Expansion", "Supports additional product categories and rule sets."),
]

econ_y = Inches(1.6)
for i, (title, desc) in enumerate(econ_benefits):
    y = econ_y + Inches(i * 0.95)

    # Number badge
    badge = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.5), y, Inches(0.3), Inches(0.25)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT_GREEN
    badge.line.fill.background()
    btf = badge.text_frame
    btf.paragraphs[0].alignment = PP_ALIGN.CENTER
    brun = btf.paragraphs[0].add_run()
    brun.text = str(i + 1)
    brun.font.size = Pt(9)
    brun.font.color.rgb = WHITE
    brun.font.bold = True
    brun.font.name = "Calibri"

    add_text_box(slide3, Inches(9.9), y - Inches(0.02), Inches(3.2), Inches(0.25),
                 title, 10, NAVY, True)
    add_text_box(slide3, Inches(9.9), y + Inches(0.23), Inches(3.2), Inches(0.45),
                 desc, 8, DARK_GRAY)

# ── BOTTOM: Impact Dimensions ──
add_text_box(slide3, Inches(0.3), Inches(6.25), Inches(12.7), Inches(0.3),
             "ILLUSTRATIVE IMPACT PRIORITIES", 11, MEDIUM_GRAY, italic=True, alignment=PP_ALIGN.CENTER)

impact_dims = [
    ("Inspection\nEfficiency", BLUE),
    ("Evidence &\nTraceability", GREEN),
    ("Inspector\nAssistance", ORANGE),
    ("Consistency", ACCENT_TEAL),
    ("Digital\nDocumentation", NAVY),
]

dim_w = Inches(2.3)
for i, (label, color) in enumerate(impact_dims):
    x = Inches(0.5) + i * (dim_w + Inches(0.2))

    dim_box = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(6.55), dim_w, Inches(0.4)
    )
    dim_box.fill.solid()
    dim_box.fill.fore_color.rgb = color
    dim_box.line.fill.background()
    dtf = dim_box.text_frame
    dtf.paragraphs[0].alignment = PP_ALIGN.CENTER
    dr = dtf.paragraphs[0].add_run()
    dr.text = label
    dr.font.size = Pt(8)
    dr.font.color.rgb = WHITE
    dr.font.bold = True
    dr.font.name = "Calibri"

add_footer(slide3, 3)


# ══════════════════════════════════════════════════════════════
# SLIDE 4: FEASIBILITY & VIABILITY + WHAT IFs
# ══════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(blank_layout)
add_sih_logo(slide4)
add_team_name_oval(slide4)
add_header_text(slide4, "FEASIBILITY & VIABILITY")

# ── TOP: 5 Feasibility Cards ──
feasibility = [
    ("01", "Technology\nFeasibility", "OCR, computer vision and AI models are technically applicable to label inspection.", BLUE),
    ("02", "Data\nFeasibility", "Product-label images and structured declarations can support development and testing.", GREEN),
    ("03", "Architecture\nFeasibility", "AI extraction and regulatory rules are separated into modular components.", ACCENT_TEAL),
    ("04", "Deployment\nFeasibility", "A web-based system can support desktop, tablet and mobile workflows.", ORANGE),
    ("05", "Scalability", "New categories and rules can be added through configurable definitions.", NAVY),
]

feas_w = Inches(2.35)
feas_y = Inches(1.15)
for i, (num, title, desc, color) in enumerate(feasibility):
    x = Inches(0.3) + i * (feas_w + Inches(0.15))

    card = add_rounded_card(slide4, x, feas_y, feas_w, Inches(1.5), CARD_BG, color)

    # Number
    badge = slide4.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Inches(0.08), feas_y + Inches(0.08),
        Inches(0.3), Inches(0.3)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    btf = badge.text_frame
    btf.paragraphs[0].alignment = PP_ALIGN.CENTER
    brun = btf.paragraphs[0].add_run()
    brun.text = num
    brun.font.size = Pt(9)
    brun.font.color.rgb = WHITE
    brun.font.bold = True
    brun.font.name = "Calibri"

    add_text_box(slide4, x + Inches(0.05), feas_y + Inches(0.4), feas_w - Inches(0.1), Inches(0.45),
                 title, 9, NAVY, True, PP_ALIGN.CENTER)
    add_text_box(slide4, x + Inches(0.05), feas_y + Inches(0.85), feas_w - Inches(0.1), Inches(0.6),
                 desc, 7.5, DARK_GRAY, alignment=PP_ALIGN.CENTER)

# ── MIDDLE LEFT: What Ifs ──
add_text_box(slide4, Inches(0.3), Inches(2.85), Inches(6.0), Inches(0.35),
             "WHAT IFs...?", 13, ACCENT_RED, True, PP_ALIGN.LEFT)

whatifs = [
    ("AI misreads text?", "Confidence scoring + evidence + inspector correction."),
    ("Image is blurry?", "Image-quality detection + adaptive recapture."),
    ("Declaration on another side?", "Multi-view package analysis."),
    ("Online info differs?", "Cross-source comparison."),
    ("Regulations change?", "Version-controlled rule engine."),
    ("AI produces false finding?", "Human-in-the-loop verification."),
]

whatif_y = Inches(3.3)
for i, (question, mitigation) in enumerate(whatifs):
    row = i // 2
    col = i % 2
    x = Inches(0.3) + col * Inches(3.1)
    y = whatif_y + row * Inches(0.85)

    card = add_rounded_card(slide4, x, y, Inches(2.95), Inches(0.75), CARD_BG, ACCENT_RED)
    add_text_box(slide4, x + Inches(0.08), y + Inches(0.02), Inches(2.8), Inches(0.22),
                 f"Q: {question}", 8, ACCENT_RED, True)
    add_text_box(slide4, x + Inches(0.08), y + Inches(0.28), Inches(2.8), Inches(0.4),
                 f"A: {mitigation}", 7, DARK_GRAY)

# ── MIDDLE RIGHT: Risk → Mitigation ──
add_text_box(slide4, Inches(6.5), Inches(2.85), Inches(6.5), Inches(0.35),
             "RISK → MITIGATION", 13, NAVY, True, PP_ALIGN.LEFT)

risks = [
    ("OCR ERROR", "CONFIDENCE + HUMAN REVIEW", BLUE),
    ("POOR IMAGE", "ADAPTIVE RECAPTURE", ACCENT_TEAL),
    ("MISSING VIEW", "MULTI-VIEW ANALYSIS", GREEN),
    ("DATA MISMATCH", "CROSS-SOURCE CHECK", ORANGE),
    ("RULE CHANGE", "VERSIONED RULE ENGINE", NAVY),
]

risk_y = Inches(3.3)
for i, (risk, mitigation, color) in enumerate(risks):
    y = risk_y + Inches(i * 0.55)

    # Risk box
    rbox = add_rounded_card(slide4, Inches(6.5), y, Inches(2.2), Inches(0.45), color)
    rtf = rbox.text_frame
    rtf.paragraphs[0].alignment = PP_ALIGN.CENTER
    rr = rtf.paragraphs[0].add_run()
    rr.text = risk
    rr.font.size = Pt(7.5)
    rr.font.color.rgb = WHITE
    rr.font.bold = True
    rr.font.name = "Calibri"
    rtf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Arrow
    add_arrow_right(slide4, Inches(8.75), y + Inches(0.1), Inches(0.3), Inches(0.18))

    # Mitigation box
    mbox = add_rounded_card(slide4, Inches(9.15), y, Inches(3.5), Inches(0.45), LIGHT_BLUE, color)
    mtf = mbox.text_frame
    mtf.paragraphs[0].alignment = PP_ALIGN.CENTER
    mr = mtf.paragraphs[0].add_run()
    mr.text = mitigation
    mr.font.size = Pt(7.5)
    mr.font.color.rgb = NAVY
    mr.font.bold = True
    mr.font.name = "Calibri"
    mtf.vertical_anchor = MSO_ANCHOR.MIDDLE

# ── BOTTOM: Before vs MetrologyAI ──
add_text_box(slide4, Inches(0.3), Inches(5.75), Inches(12.7), Inches(0.3),
             "BEFORE vs METROLOGYAI", 13, NAVY, True, PP_ALIGN.CENTER)

# BEFORE panel
before_card = add_rounded_card(slide4, Inches(0.3), Inches(6.1), Inches(5.8), Inches(0.85),
                                RGBColor(0xFF, 0xEB, 0xEE), ACCENT_RED)
add_text_box(slide4, Inches(0.5), Inches(6.12), Inches(2.0), Inches(0.25),
             "BEFORE", 10, ACCENT_RED, True)
before_items = ["Manual reading", "Manual comparison", "Fragmented evidence", "Manual documentation"]
for i, item in enumerate(before_items):
    x = Inches(0.5) + (i % 2) * Inches(2.8)
    y = Inches(6.4) + (i // 2) * Inches(0.22)
    add_text_box(slide4, x, y, Inches(2.7), Inches(0.2), f"✗ {item}", 7.5, ACCENT_RED)

# METROLOGYAI panel
after_card = add_rounded_card(slide4, Inches(6.5), Inches(6.1), Inches(6.5), Inches(0.85),
                               RGBColor(0xE8, 0xF5, 0xE9), ACCENT_GREEN)
add_text_box(slide4, Inches(6.7), Inches(6.12), Inches(3.0), Inches(0.25),
             "METROLOGYAI", 10, ACCENT_GREEN, True)
after_items = ["AI-assisted extraction", "Rule-based validation", "Evidence-linked findings", "Digital reporting"]
for i, item in enumerate(after_items):
    x = Inches(6.7) + (i % 2) * Inches(3.0)
    y = Inches(6.4) + (i // 2) * Inches(0.22)
    add_text_box(slide4, x, y, Inches(2.9), Inches(0.2), f"✓ {item}", 7.5, ACCENT_GREEN)

add_footer(slide4, 4)


# ══════════════════════════════════════════════════════════════
# SLIDE 5: PROPOSED VS EXISTING
# ══════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(blank_layout)
add_sih_logo(slide5)
add_team_name_oval(slide5)
add_header_text(slide5, "PROPOSED SOLUTION VS EXISTING APPROACHES")

# ── Comparison Table ──
add_text_box(slide5, Inches(0.3), Inches(1.15), Inches(12.7), Inches(0.35),
             "DETAILED COMPARISON", 13, NAVY, True, PP_ALIGN.CENTER)

table_data = [
    ("Dimension", "Conventional / Existing", "MetrologyAI"),
    ("Inspection", "Manual", "AI-assisted"),
    ("Package Views", "Individual views", "Multi-view"),
    ("Text Extraction", "Manual / basic OCR", "OCR + AI vision"),
    ("Image Quality", "Manual judgement", "Quality analysis"),
    ("Unreadable Text", "Often manually judged", "Confidence-aware"),
    ("Rule Checking", "Manual", "Configurable rule engine"),
    ("Evidence", "Manually documented", "Evidence-linked"),
    ("Online Comparison", "Separate", "Integrated"),
    ("Historical Comparison", "Limited", "Supported"),
    ("Risk Prioritization", "Inspector-led", "AI-assisted"),
    ("Recommendations", "Manual", "Next-best-action"),
    ("Inspector Role", "Manual processing", "Final verification"),
    ("Reporting", "Manual", "Structured digital report"),
]

table_left = Inches(0.5)
table_top = Inches(1.55)
table_w = Inches(12.3)
table_h = Inches(4.2)
num_rows = len(table_data)
num_cols = 3

table_shape = slide5.shapes.add_table(num_rows, num_cols, table_left, table_top, table_w, table_h)
table = table_shape.table

# Column widths
table.columns[0].width = Inches(3.5)
table.columns[1].width = Inches(4.4)
table.columns[2].width = Inches(4.4)

for row_idx, row_data in enumerate(table_data):
    for col_idx, cell_text in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = cell_text

        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(8.5)
            paragraph.font.name = "Calibri"

            if row_idx == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
                paragraph.alignment = PP_ALIGN.CENTER
            elif col_idx == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = NAVY
            elif col_idx == 1:
                paragraph.font.color.rgb = DARK_GRAY
            else:
                paragraph.font.color.rgb = ACCENT_GREEN
                paragraph.font.bold = True

        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Header row
        if row_idx == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
        elif row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BLUE
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE

# ── BOTTOM: Three Feature Pillars ──
add_text_box(slide5, Inches(0.3), Inches(5.9), Inches(12.7), Inches(0.3),
             "FROM OCR TO INSPECTION INTELLIGENCE", 14, NAVY, True, PP_ALIGN.CENTER)

pillars = [
    ("SEE", "Multi-View Intelligence", "Analyze front, back, side, top and bottom views of packaged commodities.", BLUE),
    ("UNDERSTAND", "Evidence + Rules + Context", "AI-driven extraction validated against configurable compliance rules with full evidence chains.", ORANGE),
    ("VERIFY", "Human-in-the-Loop", "Inspector corrections trigger smart revalidation, ensuring final accuracy and accountability.", GREEN),
]

pillar_w = Inches(3.9)
for i, (title, subtitle, desc, color) in enumerate(pillars):
    x = Inches(0.5) + i * (pillar_w + Inches(0.25))

    card = add_rounded_card(slide5, x, Inches(6.25), pillar_w, Inches(0.75), color)

    add_text_box(slide5, x + Inches(0.1), Inches(6.27), Inches(0.8), Inches(0.25),
                 title, 14, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide5, x + Inches(0.9), Inches(6.27), Inches(2.8), Inches(0.2),
                 subtitle, 9, WHITE, True)
    add_text_box(slide5, x + Inches(0.1), Inches(6.55), pillar_w - Inches(0.2), Inches(0.4),
                 desc, 7, RGBColor(0xE0, 0xE0, 0xE0), alignment=PP_ALIGN.LEFT)

add_footer(slide5, 5)


# ══════════════════════════════════════════════════════════════
# SLIDE 6: RESEARCH & REFERENCES
# ══════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(blank_layout)
add_sih_logo(slide6)
add_team_name_oval(slide6)
add_header_text(slide6, "RESEARCH & REFERENCES")

# ── LEFT: Legal & Regulatory ──
add_text_box(slide6, Inches(0.3), Inches(1.15), Inches(4.0), Inches(0.35),
             "LEGAL & REGULATORY REFERENCES", 13, NAVY, True, PP_ALIGN.LEFT)

legal_refs = [
    ("01", "Legal Metrology Act, 2009", "Central legislation governing metrology standards in India."),
    ("02", "Legal Metrology (Packaged Commodities) Rules, 2011", "Rules for compliance of packaged commodity labels."),
    ("03", "Department of Consumer Affairs", "Nodal department for Legal Metrology administration."),
    ("04", "Official Notifications & Amendments", "Relevant notifications under the Legal Metrology framework."),
]

legal_y = Inches(1.65)
for i, (num, title, desc) in enumerate(legal_refs):
    y = legal_y + Inches(i * 1.1)

    badge = slide6.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.3), y, Inches(0.35), Inches(0.35)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = BLUE
    badge.line.fill.background()
    btf = badge.text_frame
    btf.paragraphs[0].alignment = PP_ALIGN.CENTER
    brun = btf.paragraphs[0].add_run()
    brun.text = num
    brun.font.size = Pt(9)
    brun.font.color.rgb = WHITE
    brun.font.bold = True
    brun.font.name = "Calibri"

    add_text_box(slide6, Inches(0.8), y, Inches(3.5), Inches(0.25), title, 9, NAVY, True)
    add_text_box(slide6, Inches(0.8), y + Inches(0.28), Inches(3.5), Inches(0.55), desc, 7.5, DARK_GRAY)

# ── CENTER: Research Areas ──
add_text_box(slide6, Inches(4.5), Inches(1.15), Inches(4.0), Inches(0.35),
             "RESEARCH AREAS", 13, NAVY, True, PP_ALIGN.LEFT)

research = [
    ("Computer Vision", "OCR and scene-text recognition for label extraction.", BLUE),
    ("Multimodal AI", "Image + text understanding for holistic label analysis.", GREEN),
    ("Explainable AI", "Confidence and evidence-based decisions for transparency.", ORANGE),
    ("Human-AI Interaction", "Human verification in AI-assisted inspection workflows.", NAVY),
]

research_y = Inches(1.65)
for i, (title, desc, color) in enumerate(research):
    y = research_y + Inches(i * 1.1)

    card = add_rounded_card(slide6, Inches(4.5), y, Inches(4.2), Inches(0.95), CARD_BG, color)

    # Color bar
    bar = slide6.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4.5), y, Inches(0.1), Inches(0.95)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text_box(slide6, Inches(4.75), y + Inches(0.08), Inches(3.8), Inches(0.25),
                 title, 10, NAVY, True)
    add_text_box(slide6, Inches(4.75), y + Inches(0.38), Inches(3.8), Inches(0.5),
                 desc, 8, DARK_GRAY)

# ── RIGHT: Technology References ──
add_text_box(slide6, Inches(9.0), Inches(1.15), Inches(4.0), Inches(0.35),
             "TECHNOLOGY REFERENCES", 13, NAVY, True, PP_ALIGN.LEFT)

tech_refs = [
    ("React", "Component-based UI library"),
    ("TypeScript", "Type-safe JavaScript"),
    ("Convex", "Serverless backend & database"),
    ("Tailwind CSS", "Utility-first CSS framework"),
    ("Gemini AI", "Vision + language model"),
    ("Vite", "Build tool & dev server"),
    ("Recharts", "Charting library"),
]

tech_y = Inches(1.65)
for i, (tech, desc) in enumerate(tech_refs):
    y = tech_y + Inches(i * 0.6)

    colors = [BLUE, ACCENT_TEAL, GREEN, ORANGE, NAVY, ACCENT_BLUE, ACCENT_GREEN]
    dot = slide6.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(9.0), y + Inches(0.05), Inches(0.18), Inches(0.18)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = colors[i]
    dot.line.fill.background()

    add_text_box(slide6, Inches(9.3), y, Inches(1.5), Inches(0.22), tech, 9, NAVY, True)
    add_text_box(slide6, Inches(10.8), y, Inches(2.5), Inches(0.22), desc, 8, DARK_GRAY)

# ── BOTTOM: References ──
add_text_box(slide6, Inches(0.3), Inches(6.0), Inches(12.7), Inches(0.35),
             "REFERENCES", 13, NAVY, True, PP_ALIGN.CENTER)

refs_box = add_rounded_card(slide6, Inches(0.3), Inches(6.35), Inches(12.7), Inches(0.55),
                             CARD_BG, NAVY)

refs = [
    "[ADD VERIFIED REFERENCE — Legal Metrology Act, 2009]",
    "[ADD VERIFIED REFERENCE — Legal Metrology (Packaged Commodities) Rules, 2011]",
    "[ADD VERIFIED REFERENCE — Department of Consumer Affairs, Government of India]",
    "[ADD VERIFIED REFERENCE — Relevant court judgments or regulatory orders]",
]

for i, ref in enumerate(refs):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + col * Inches(6.2)
    y = Inches(6.4) + row * Inches(0.2)
    add_text_box(slide6, x, y, Inches(6.0), Inches(0.2), ref, 7, MEDIUM_GRAY, italic=True)

add_footer(slide6, 6)


# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
output_path = "/home/daytona/codebase/MetrologyAI_SIH_2026_Final.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Slides: {len(prs.slides)}")
print(f"   Format: Widescreen 16:9")
