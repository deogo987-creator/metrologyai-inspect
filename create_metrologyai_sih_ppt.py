#!/usr/bin/env python3
"""
MetrologyAI SIH 2026 Presentation Generator
Edits the uploaded SIH template with MetrologyAI content.
Follows Maitri AI winner's structure with diagrams and visual elements.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# Colors matching SIH template theme
NAVY = RGBColor(0x00, 0x2B, 0x5C)
BLUE = RGBColor(0x00, 0x70, 0xC0)
LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x00, 0xA6, 0x5A)
RED = RGBColor(0xCC, 0x00, 0x00)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
LIGHT_RED = RGBColor(0xFF, 0xEB, 0xEE)
LIGHT_ORANGE = RGBColor(0xFF, 0xF3, 0xE0)
SAFFRON = RGBColor(0xFF, 0x99, 0x33)


def add_text_box(slide, left, top, width, height, text, font_size=12, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def add_rich_text_box(slide, left, top, width, height, lines, default_size=11, default_color=DARK_GRAY):
    """Add a text box with multiple formatted lines."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if isinstance(line, dict):
            p.text = line.get('text', '')
            p.font.size = Pt(line.get('size', default_size))
            p.font.bold = line.get('bold', False)
            p.font.color.rgb = line.get('color', default_color)
            p.alignment = line.get('align', PP_ALIGN.LEFT)
            p.space_after = Pt(line.get('space_after', 2))
            p.space_before = Pt(line.get('space_before', 0))
        else:
            p.text = str(line)
            p.font.size = Pt(default_size)
            p.font.color.rgb = default_color
            p.space_after = Pt(2)
    
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color=LIGHT_BLUE, border_color=BLUE, text="", font_size=10, font_color=DARK_GRAY):
    """Add a rounded rectangle with text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
    
    return shape


def add_arrow(slide, left, top, width, height, color=BLUE):
    """Add a right arrow."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_down_arrow(slide, left, top, width, height, color=BLUE):
    """Add a down arrow."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, fill_color=BLUE, text="", font_size=10, font_color=WHITE):
    """Add a circle with text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top), Inches(size), Inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    
    return shape


def update_text_in_shape(shape, new_text):
    """Update all text in a shape."""
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = new_text
        # Also set the first paragraph directly
        if shape.text_frame.paragraphs:
            shape.text_frame.paragraphs[0].text = new_text


def set_team_name_oval(slide, team_name="Zillion\nMindsSashakt"):
    """Update the 'Your Team Name' oval on content slides."""
    for shape in slide.shapes:
        if shape.name.startswith("Oval") and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = team_name
                para.alignment = PP_ALIGN.CENTER


def edit_slide_1(prs):
    """Slide 1: Title Page — PS26034, MetrologyAI, Team info."""
    slide = prs.slides[0]
    
    for shape in slide.shapes:
        if shape.name == "TextBox 9" and shape.has_text_frame:
            # Replace the bullet points with MetrologyAI content
            lines = [
                "Problem Statement ID – PS26034",
                "Problem Statement Title – AI-Assisted Legal Metrology",
                "Compliance Inspection System",
                "Theme – Software",
                "PS Category – Software",
                "Team ID – 57893",
                "Team Name – Zillion MindsSashakt",
            ]
            for i, para in enumerate(shape.text_frame.paragraphs):
                if i < len(lines):
                    for run in para.runs:
                        run.text = lines[i]
                    if not para.runs:
                        para.text = lines[i]
        
        if shape.name == "Subtitle 3" and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = "TITLE PAGE"
    
    print("Slide 1: Title Page updated")


def edit_slide_2(prs):
    """Slide 2: Problem + Solution + Unique Solutions + Prototype (Maitri-style multi-column)."""
    slide = prs.slides[1]
    
    # Update title
    for shape in slide.shapes:
        if shape.name == "Title 1" and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = "METROLOGYAI"
        
        # Update team name oval
        if shape.name.startswith("Oval") and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = "Zillion\nMindsSashakt"
                para.alignment = PP_ALIGN.CENTER
    
    # Remove old text box content and replace with structured content
    for shape in slide.shapes:
        if shape.name == "TextBox 8" and shape.has_text_frame:
            # Clear existing text
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ""
            
            # Build structured content (Maitri-style)
            lines = [
                {'text': 'Understanding The Problem', 'size': 16, 'bold': True, 'color': NAVY, 'space_after': 6},
                {'text': '• Manual Inspection — Inspectors manually examine labels, making checks time-consuming', 'size': 10, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Multi-Surface Declarations — Required info appears on different sides of the package', 'size': 10, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Real-World Image Issues — Blur, glare, distortion reduce readability', 'size': 10, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Uncertain AI Results — System confuses unreadable with genuinely missing', 'size': 10, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Fragmented Verification — Physical, online, historical data checked separately', 'size': 10, 'color': DARK_GRAY, 'space_after': 8},
                {'text': 'Our Solution — MetrologyAI', 'size': 16, 'bold': True, 'color': NAVY, 'space_after': 6},
                {'text': '• Takes product images → AI reads label → Checks 11 Legal Metrology rules', 'size': 10, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Gives compliance score (0-100) with evidence chain for every finding', 'size': 10, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Inspector reviews, corrects, and makes final decision', 'size': 10, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Digital PDF report with complete audit trail', 'size': 10, 'color': DARK_GRAY, 'space_after': 8},
                {'text': 'Our Unique Solutions', 'size': 16, 'bold': True, 'color': NAVY, 'space_after': 6},
                {'text': '• Multi-View Intelligence — Analyzes front, back, side, top, bottom views', 'size': 10, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Missing ≠ Unreadable — Separates absent from low-confidence detection', 'size': 10, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Adaptive Re-Capture — Identifies poor quality, recommends improvement', 'size': 10, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Evidence Chain — Links finding to image, OCR, field, and rule', 'size': 10, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Risk Prioritization — AI scores every inspection for risk level', 'size': 10, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Human-in-the-Loop — Corrections trigger automatic revalidation', 'size': 10, 'color': DARK_GRAY, 'space_after': 2},
            ]
            
            for i, line in enumerate(lines):
                if i == 0:
                    p = shape.text_frame.paragraphs[0]
                else:
                    p = shape.text_frame.add_paragraph()
                p.text = line['text']
                p.font.size = Pt(line['size'])
                p.font.bold = line.get('bold', False)
                p.font.color.rgb = line['color']
                p.space_after = Pt(line.get('space_after', 2))
    
    # Add USPs on the right side (Maitri-style)
    usps = [
        ("Real-time\ncompliance\nscoring", GREEN),
        ("Evidence-backed\nfindings for\nevery violation", BLUE),
        ("Multi-view\npackage\nanalysis", SAFFRON),
        ("Inspector\ncorrections with\nauto-revalidation", NAVY),
    ]
    
    for i, (text, color) in enumerate(usps):
        x = 9.8 + (i % 2) * 1.7
        y = 1.5 + (i // 2) * 2.2
        shape = add_rounded_rect(slide, x, y, 1.5, 1.8, LIGHT_BLUE, color, text, 8, DARK_GRAY)
    
    # Add USPs label
    add_text_box(slide, 9.8, 1.1, 3.2, 0.4, "USPs {", 14, True, NAVY)
    
    print("Slide 2: Problem + Solution + Unique Solutions updated")


def edit_slide_3(prs):
    """Slide 3: Technical Approach — Methodology + Architecture + Tech Stack."""
    slide = prs.slides[2]
    
    # Update title
    for shape in slide.shapes:
        if shape.name == "Title 1" and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = "TECHNICAL APPROACH"
        
        if shape.name.startswith("Oval") and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = "Zillion\nMindsSashakt"
                para.alignment = PP_ALIGN.CENTER
    
    # Clear old text and replace with structured content
    for shape in slide.shapes:
        if shape.name == "TextBox 8" and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ""
            
            lines = [
                {'text': 'Methodology', 'size': 14, 'bold': True, 'color': NAVY, 'space_after': 4},
                {'text': '01 CAPTURE → 02 QUALITY CHECK → 03 AI EXTRACTION → 04 STRUCTURE → 05 VALIDATE → 06 EVIDENCE → 07 RISK → 08 REVIEW → 09 REPORT', 'size': 9, 'color': DARK_GRAY, 'space_after': 6},
                {'text': 'Technical Approach', 'size': 14, 'bold': True, 'color': NAVY, 'space_after': 4},
                {'text': '• Image Acquisition — Capture multi-view product images (front, back, left, right, top, bottom)', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Image Quality Analysis — Client-side canvas: blur, brightness, glare, resolution, perspective', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• AI Vision + OCR — Gemini Vision API extracts 11 fields with confidence scores', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Rule Validation — 11 Legal Metrology rules checked against extracted fields', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Evidence Engine — Every finding linked to image region, OCR text, field, and rule', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Risk Scoring — Compliance score, missing declarations, low confidence, anomalies', 'size': 9, 'color': DARK_GRAY, 'space_after': 6},
                {'text': 'Tech Stack', 'size': 14, 'bold': True, 'color': NAVY, 'space_after': 4},
                {'text': 'Frontend: React + TypeScript + Vite + Tailwind CSS + shadcn/ui', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': 'Backend: Convex (Serverless) | Database: Convex DB (Real-time)', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': 'AI/Vision: Google Gemini Vision API (Multi-model fallback) | Auth: Convex Auth', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': 'Analytics: Recharts | Deployment: Freebuff Platform', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
            ]
            
            for i, line in enumerate(lines):
                if i == 0:
                    p = shape.text_frame.paragraphs[0]
                else:
                    p = shape.text_frame.add_paragraph()
                p.text = line['text']
                p.font.size = Pt(line['size'])
                p.font.bold = line.get('bold', False)
                p.font.color.rgb = line['color']
                p.space_after = Pt(line.get('space_after', 2))
    
    # Add flowchart on the right side
    flow_steps = [
        ("IMAGE", LIGHT_BLUE),
        ("QUALITY\nCHECK", LIGHT_GREEN),
        ("AI VISION\n+ OCR", LIGHT_ORANGE),
        ("RULE\nVALIDATION", LIGHT_RED),
        ("SCORE +\nEVIDENCE", LIGHT_BLUE),
        ("INSPECTOR\nREVIEW", LIGHT_GREEN),
        ("PDF\nREPORT", LIGHT_ORANGE),
    ]
    
    for i, (text, color) in enumerate(flow_steps):
        y = 1.3 + i * 0.75
        add_rounded_rect(slide, 10.5, y, 2.2, 0.6, color, NAVY, text, 8, DARK_GRAY)
        if i < len(flow_steps) - 1:
            add_down_arrow(slide, 11.4, y + 0.58, 0.3, 0.15, NAVY)
    
    print("Slide 3: Technical Approach updated")


def edit_slide_4(prs):
    """Slide 4: Feasibility + What Ifs + Challenge Mitigation."""
    slide = prs.slides[3]
    
    # Update title
    for shape in slide.shapes:
        if shape.name == "Title 1" and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = "FEASIBILITY AND VIABILITY"
        
        if shape.name.startswith("Oval") and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = "Zillion\nMindsSashakt"
                para.alignment = PP_ALIGN.CENTER
    
    # Clear old text and replace
    for shape in slide.shapes:
        if shape.name == "TextBox 8" and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ""
            
            lines = [
                {'text': 'Feasibility Analysis', 'size': 14, 'bold': True, 'color': NAVY, 'space_after': 4},
                {'text': '01 Technology — Gemini Vision API is production-ready; React + Convex are proven frameworks', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '02 Data — Product labels easy to collect; 11 rules publicly documented in Legal Metrology Rules 2011', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '03 Architecture — AI extraction and rule validation are modular, independently testable', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '04 Deployment — Web-based; works on desktop, tablet, mobile; no app install required', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '05 Scalability — New categories added by updating RULE_REQUIREMENTS config', 'size': 9, 'color': DARK_GRAY, 'space_after': 6},
                {'text': 'What Ifs...? — Challenge → Mitigation', 'size': 14, 'bold': True, 'color': NAVY, 'space_after': 4},
                {'text': '• AI misreads text → Confidence scoring + evidence chain + inspector correction', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Image is blurry → Quality detection + adaptive recapture recommendations', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Declaration on another side → Multi-view analysis (up to 6 views)', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Regulations change → Version-controlled rule engine', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• AI produces false finding → Human-in-the-loop verification', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• API fails → Multi-model fallback (4 Gemini models in sequence)', 'size': 9, 'color': DARK_GRAY, 'space_after': 6},
                {'text': 'Before MetrologyAI → After MetrologyAI', 'size': 12, 'bold': True, 'color': NAVY, 'space_after': 4},
                {'text': 'BEFORE: Manual reading | Manual comparison | Fragmented evidence | Paper reports', 'size': 9, 'color': RED, 'space_after': 2},
                {'text': 'AFTER: AI-assisted extraction | Rule-based validation | Evidence-linked | Digital reports', 'size': 9, 'color': GREEN, 'space_after': 2},
            ]
            
            for i, line in enumerate(lines):
                if i == 0:
                    p = shape.text_frame.paragraphs[0]
                else:
                    p = shape.text_frame.add_paragraph()
                p.text = line['text']
                p.font.size = Pt(line['size'])
                p.font.bold = line.get('bold', False)
                p.font.color.rgb = line['color']
                p.space_after = Pt(line.get('space_after', 2))
    
    print("Slide 4: Feasibility updated")


def edit_slide_5(prs):
    """Slide 5: Impact and Benefits."""
    slide = prs.slides[4]
    
    # Update title
    for shape in slide.shapes:
        if shape.name == "Title 1" and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = "IMPACT AND BENEFITS"
        
        if shape.name.startswith("Oval") and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = "Zillion\nMindsSashakt"
                para.alignment = PP_ALIGN.CENTER
    
    # Clear old text and replace
    for shape in slide.shapes:
        if shape.name == "TextBox 8" and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ""
            
            lines = [
                {'text': 'Inspector Benefits', 'size': 14, 'bold': True, 'color': NAVY, 'space_after': 4},
                {'text': '• Faster inspection — AI reads 11 fields in seconds, not minutes', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Better accuracy — Confidence scores show what is reliable', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Complete evidence — Every finding linked to source image and rule', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Risk prioritization — Focus on high-risk products first', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Smart revalidation — Corrections automatically update compliance status', 'size': 9, 'color': DARK_GRAY, 'space_after': 6},
                {'text': 'Operational Benefits', 'size': 14, 'bold': True, 'color': NAVY, 'space_after': 4},
                {'text': '• Consistent checking — Same 11 rules applied every time', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Reduced false positives — Missing vs unreadable distinction', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Scalable — Supports more inspections without proportional staff increase', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Data-driven — Analytics show patterns across inspections', 'size': 9, 'color': DARK_GRAY, 'space_after': 6},
                {'text': 'Social & Economic Benefits', 'size': 14, 'bold': True, 'color': NAVY, 'space_after': 4},
                {'text': '• Better consumer protection through consistent enforcement', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Faster detection of non-compliant products', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Digital evidence supports legal proceedings', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Supports government Digital India initiative', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '• Less manual effort, fewer re-inspections, digital documentation', 'size': 9, 'color': DARK_GRAY, 'space_after': 2},
            ]
            
            for i, line in enumerate(lines):
                if i == 0:
                    p = shape.text_frame.paragraphs[0]
                else:
                    p = shape.text_frame.add_paragraph()
                p.text = line['text']
                p.font.size = Pt(line['size'])
                p.font.bold = line.get('bold', False)
                p.font.color.rgb = line['color']
                p.space_after = Pt(line.get('space_after', 2))
    
    # Add impact dimensions visual on the right
    dimensions = [
        ("Inspection\nEfficiency", "35", GREEN),
        ("Evidence &\nTraceability", "25", BLUE),
        ("Inspector\nAssistance", "20", SAFFRON),
        ("Consistency", "10", NAVY),
        ("Digital\nDocumentation", "10", GRAY),
    ]
    
    for i, (label, pct, color) in enumerate(dimensions):
        y = 1.3 + i * 1.15
        # Bar
        bar_width = float(pct) / 35 * 2.5
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(10.0), Inches(y), Inches(bar_width), Inches(0.4)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = f"{pct}%"
        p.font.size = Pt(8)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        add_text_box(slide, 10.0, y + 0.42, 2.5, 0.5, label, 7, True, DARK_GRAY)
    
    # Chart title
    add_text_box(slide, 10.0, 0.9, 3.0, 0.4, "Impact Priorities", 10, True, NAVY, PP_ALIGN.CENTER)
    
    print("Slide 5: Impact and Benefits updated")


def edit_slide_6(prs):
    """Slide 6: Research and References."""
    slide = prs.slides[5]
    
    # Update title
    for shape in slide.shapes:
        if shape.name == "Title 1" and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = "RESEARCH AND REFERENCES"
        
        if shape.name.startswith("Oval") and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = "Zillion\nMindsSashakt"
                para.alignment = PP_ALIGN.CENTER
    
    # Clear old text and replace
    for shape in slide.shapes:
        if shape.name == "TextBox 8" and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ""
            
            lines = [
                {'text': 'Legal & Regulatory References', 'size': 14, 'bold': True, 'color': NAVY, 'space_after': 4},
                {'text': '01 Legal Metrology Act, 2009 — Government of India', 'size': 10, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '02 Legal Metrology (Packaged Commodities) Rules, 2011', 'size': 10, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '03 Department of Consumer Affairs — Legal Metrology Division', 'size': 10, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '04 Official Gazette Notifications and Amendments', 'size': 10, 'color': DARK_GRAY, 'space_after': 6},
                {'text': 'Research Areas', 'size': 14, 'bold': True, 'color': NAVY, 'space_after': 4},
                {'text': '01 Computer Vision & OCR — Scene-text recognition for label analysis', 'size': 10, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '02 Multimodal AI — Vision-language models for regulatory compliance', 'size': 10, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '03 Explainable AI — Confidence scoring and evidence-based decisions', 'size': 10, 'color': DARK_GRAY, 'space_after': 2},
                {'text': '04 Human-AI Interaction — Inspector-in-the-loop verification patterns', 'size': 10, 'color': DARK_GRAY, 'space_after': 6},
                {'text': 'Technology References', 'size': 14, 'bold': True, 'color': NAVY, 'space_after': 4},
                {'text': 'React (react.dev) | TypeScript (typescriptlang.org) | Vite (vitejs.dev)', 'size': 10, 'color': DARK_GRAY, 'space_after': 2},
                {'text': 'Tailwind CSS (tailwindcss.com) | Convex (convex.dev) | shadcn/ui (ui.shadcn.com)', 'size': 10, 'color': DARK_GRAY, 'space_after': 2},
                {'text': 'Google Gemini Vision API (ai.google.dev) | Recharts (recharts.org)', 'size': 10, 'color': DARK_GRAY, 'space_after': 2},
                {'text': 'Lucide Icons (lucide.dev) | Sonner Toast (sonner.emilkowal.ski)', 'size': 10, 'color': DARK_GRAY, 'space_after': 6},
                {'text': '[ADD VERIFIED RESEARCH PAPERS OR ACADEMIC REFERENCES HERE]', 'size': 10, 'bold': True, 'color': ORANGE, 'space_after': 2},
            ]
            
            for i, line in enumerate(lines):
                if i == 0:
                    p = shape.text_frame.paragraphs[0]
                else:
                    p = shape.text_frame.add_paragraph()
                p.text = line['text']
                p.font.size = Pt(line['size'])
                p.font.bold = line.get('bold', False)
                p.font.color.rgb = line['color']
                p.space_after = Pt(line.get('space_after', 2))
    
    print("Slide 6: Research and References updated")


def main():
    print("Opening SIH template...")
    prs = Presentation('SIH2026-IDEA-Presentation-Format.pptx')
    
    print(f"Template: {len(prs.slides)} slides, {prs.slide_width/914400:.1f}x{prs.slide_height/914400:.1f} inches")
    
    # Edit each slide
    edit_slide_1(prs)  # Title Page
    edit_slide_2(prs)  # Problem + Solution + Unique Solutions
    edit_slide_3(prs)  # Technical Approach
    edit_slide_4(prs)  # Feasibility + What Ifs
    edit_slide_5(prs)  # Impact and Benefits
    edit_slide_6(prs)  # Research and References
    
    # Save
    output_path = 'MetrologyAI_SIH_2026_Final.pptx'
    prs.save(output_path)
    print(f"\nSaved: {output_path}")
    
    # Verify
    prs2 = Presentation(output_path)
    print(f"Verification: {len(prs2.slides)} slides")
    for i, slide in enumerate(prs2.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()[:60]
                if t:
                    texts.append(t)
        print(f"  Slide {i+1}: {texts[0] if texts else 'empty'}")


if __name__ == '__main__':
    main()
