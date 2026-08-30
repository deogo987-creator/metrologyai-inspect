#!/usr/bin/env python3
"""
MetrologyAI - SIH 2026 Presentation Generator
Matches the EXACT template layout from uploaded screenshots.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors matching SIH template ──
NAVY = RGBColor(0x1B, 0x3A, 0x6B)        # Dark navy for "SMART INDIA HACKATHON 2026"
GREEN = RGBColor(0x27, 0xAE, 0x60)       # SIH green for "HACKATHON 2026"
ORANGE = RGBColor(0xE8, 0x6C, 0x00)      # SIH orange
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)       # Pure black for titles
DARK_BLUE = RGBColor(0x00, 0x3E, 0x7E)   # Blue for subtitles
FOOTER_BLUE = RGBColor(0x00, 0x56, 0xA0) # Footer bar blue
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)  # Hexagon decorative

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]


def add_footer(slide, page_num):
    """Add blue footer bar exactly matching template."""
    footer = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.05),
        Inches(13.333), Inches(0.45)
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = FOOTER_BLUE
    footer.line.fill.background()

    # Footer text
    tf = footer.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "@SIH Idea submission- Template"
    run.font.size = Pt(11)
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"

    # Page number in bottom right
    num_box = slide.shapes.add_textbox(
        Inches(12.5), Inches(7.1),
        Inches(0.5), Inches(0.35)
    )
    tf2 = num_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run2 = p2.add_run()
    run2.text = str(page_num)
    run2.font.size = Pt(14)
    run2.font.color.rgb = WHITE
    run2.font.bold = True
    run2.font.name = "Calibri"


def add_team_name_oval(slide):
    """Add 'Your Team Name' oval in top-left corner matching template."""
    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(0.3), Inches(0.15),
        Inches(1.3), Inches(0.85)
    )
    oval.fill.background()
    oval.line.color.rgb = RGBColor(0x44, 0x44, 0x44)  # Dark gray border
    oval.line.width = Pt(1.5)

    tf = oval.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Your\nTeam Name"
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
    run.font.bold = False
    run.font.name = "Calibri"


def add_sih_logo_top_right(slide):
    """Add SIH logo and text in top-right corner matching template."""
    # "SMART INDIA" text
    smart_box = slide.shapes.add_textbox(
        Inches(11.3), Inches(0.05),
        Inches(1.9), Inches(0.35)
    )
    tf = smart_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "SMART INDIA"
    run.font.size = Pt(12)
    run.font.color.rgb = NAVY
    run.font.bold = True
    run.font.name = "Calibri"

    # "HACKATHON 2026" text
    hack_box = slide.shapes.add_textbox(
        Inches(11.3), Inches(0.35),
        Inches(1.9), Inches(0.35)
    )
    tf2 = hack_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run2 = p2.add_run()
    run2.text = "HACKATHON 2026"
    run2.font.size = Pt(12)
    run2.font.color.rgb = GREEN
    run2.font.bold = True
    run2.font.name = "Calibri"

    # SIH badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(12.1), Inches(0.72),
        Inches(0.5), Inches(0.3)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = ORANGE
    badge.line.fill.background()
    btf = badge.text_frame
    btf.paragraphs[0].alignment = PP_ALIGN.CENTER
    brun = btf.paragraphs[0].add_run()
    brun.text = "SIH"
    brun.font.size = Pt(8)
    brun.font.color.rgb = WHITE
    brun.font.bold = True
    brun.font.name = "Calibri"


def add_title(slide, title_text, y=Inches(0.3)):
    """Add centered title text matching template style."""
    title_box = slide.shapes.add_textbox(
        Inches(1.5), y,
        Inches(10), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(32)
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = "Calibri"


def add_bullet_points(slide, bullets, start_y=Inches(1.8), font_size=Pt(18)):
    """Add bullet points matching template format."""
    tb = slide.shapes.add_textbox(
        Inches(1.0), start_y,
        Inches(11), Inches(5)
    )
    tf = tb.text_frame
    tf.word_wrap = True

    for i, bullet_text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.space_after = Pt(8)
        p.space_before = Pt(4)

        run = p.add_run()
        run.text = "•  " + bullet_text
        run.font.size = font_size
        run.font.color.rgb = BLACK
        run.font.name = "Calibri"


def add_blue_underlined_heading(slide, text, y=Inches(1.8)):
    """Add blue underlined heading like 'Proposed Solution' on slide 2."""
    tb = slide.shapes.add_textbox(
        Inches(1.0), y,
        Inches(11), Inches(0.6)
    )
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(22)
    run.font.color.rgb = DARK_BLUE
    run.font.bold = False
    run.font.underline = True
    run.font.name = "Calibri"


def add_hexagon_decorations(slide):
    """Add decorative hexagon shapes on the right side like slide 1."""
    # Large hexagon (background)
    hex1 = slide.shapes.add_shape(
        MSO_SHAPE.HEXAGON,
        Inches(7.5), Inches(1.5),
        Inches(5), Inches(5)
    )
    hex1.fill.solid()
    hex1.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    hex1.line.fill.background()
    hex1.rotation = 15

    # Medium hexagon
    hex2 = slide.shapes.add_shape(
        MSO_SHAPE.HEXAGON,
        Inches(8.5), Inches(2.5),
        Inches(3.5), Inches(3.5)
    )
    hex2.fill.solid()
    hex2.fill.fore_color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
    hex2.line.fill.background()
    hex2.rotation = 30

    # Small hexagon
    hex3 = slide.shapes.add_shape(
        MSO_SHAPE.HEXAGON,
        Inches(7.0), Inches(3.0),
        Inches(2), Inches(2)
    )
    hex3.fill.solid()
    hex3.fill.fore_color.rgb = RGBColor(0xC0, 0xC0, 0xC0)
    hex3.line.fill.background()
    hex3.rotation = -10


def add_sih_brain_logo_large(slide):
    """Add large SIH brain logo on slide 1."""
    # Brain circle
    brain = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(9.0), Inches(2.0),
        Inches(3.5), Inches(3.5)
    )
    brain.fill.solid()
    brain.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)
    brain.line.fill.background()

    # "SIH" text in center of brain
    sih_text = slide.shapes.add_textbox(
        Inches(9.8), Inches(3.8),
        Inches(2), Inches(1)
    )
    stf = sih_text.text_frame
    sp = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sr = sp.add_run()
    sr.text = "SIH"
    sr.font.size = Pt(28)
    sr.font.color.rgb = NAVY
    sr.font.bold = True
    sr.font.name = "Calibri"


# ══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE PAGE (matches template exactly)
# ══════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(blank_layout)

# "SMART INDIA HACKATHON 2026" header - large, centered, blue
hdr = slide1.shapes.add_textbox(
    Inches(1.5), Inches(0.2),
    Inches(10), Inches(0.6)
)
htf = hdr.text_frame
hp = htf.paragraphs[0]
hp.alignment = PP_ALIGN.CENTER
hr = hp.add_run()
hr.text = "SMART INDIA HACKATHON 2026"
hr.font.size = Pt(36)
hr.font.color.rgb = NAVY
hr.font.bold = True
hr.font.name = "Calibri"

# "TITLE PAGE" subtitle
sub = slide1.shapes.add_textbox(
    Inches(3), Inches(1.0),
    Inches(7), Inches(0.5)
)
stf = sub.text_frame
sp = stf.paragraphs[0]
sp.alignment = PP_ALIGN.CENTER
sr = sp.add_run()
sr.text = "TITLE PAGE"
sr.font.size = Pt(24)
sr.font.color.rgb = BLACK
sr.font.bold = True
sr.font.name = "Calibri"

# Bullet points on left (matching template exactly)
title_bullets = [
    "Problem Statement ID – PS26034",
    "Problem Statement Title – AI-Assisted Legal Metrology Label Compliance Inspection System",
    "Theme – Software",
    "PS Category – Software",
    "Team ID – [TEAM ID]",
    "Team Name – [TEAM NAME] (Registered on portal)"
]

add_bullet_points(slide1, title_bullets, start_y=Inches(2.0), font_size=Pt(20))

# Hexagon decorations on right (matching template)
add_hexagon_decorations(slide1)

# Large SIH brain logo on right (matching template)
add_sih_brain_logo_large(slide1)

# SIH logo top-right corner
add_sih_logo_top_right(slide1)


# ══════════════════════════════════════════════════════════════
# SLIDE 2: IDEA TITLE (matches template exactly)
# ══════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(blank_layout)

add_team_name_oval(slide2)
add_sih_logo_top_right(slide2)
add_title(slide2, "METROLOGYAI")

# Blue underlined heading
add_blue_underlined_heading(slide2, "Proposed Solution — AI-Assisted Legal Metrology Compliance Inspection System")

# Bullet points
bullets_slide2 = [
    "MetrologyAI is a software system that assists inspectors in checking packaged commodity labels under the Legal Metrology (Packaged Commodities) Rules, 2011.",
    "The system uses product images, OCR, computer vision and a configurable compliance rule engine to automate label analysis.",
    "AI assists with: image analysis, text extraction, declaration identification, compliance validation, evidence generation, risk prioritization, cross-source comparison, inspection recommendations and digital reporting.",
    "The inspector remains responsible for the final verification and decision — the system does not replace the authorized inspector.",
    "Unique capabilities: Multi-View Intelligence (front, back, side, top, bottom), Missing ≠ Unreadable distinction, Adaptive Re-Capture, Evidence Chain linking, Cross-Source Verification and Human-in-the-Loop revalidation."
]

add_bullet_points(slide2, bullets_slide2, start_y=Inches(2.8), font_size=Pt(18))

add_footer(slide2, 2)


# ══════════════════════════════════════════════════════════════
# SLIDE 3: TECHNICAL APPROACH (matches template exactly)
# ══════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(blank_layout)

add_team_name_oval(slide3)
add_sih_logo_top_right(slide3)
add_title(slide3, "TECHNICAL APPROACH")

bullets_slide3 = [
    "Technologies used: React, TypeScript, Vite, Tailwind CSS (Frontend) • Convex Serverless Backend & Database • Gemini AI Vision + OCR (AI/ML) • Recharts (Analytics)",
    "Methodology: Multi-view product image capture → Image quality analysis (blur, glare, visibility) → OCR + AI vision text extraction → Structured declaration extraction (MRP, manufacturer, net quantity, consumer care, country of origin) → Compliance rule validation → Evidence engine linking findings to source → Risk prioritization → Inspector review → Digital report generation",
    "Implementation flow: Dashboard → New Inspection → Upload product images (multi-view) → AI analysis → Compliance results with evidence chains → Inspector corrections → Smart revalidation → Export report",
    "Architecture: Modular design separating AI extraction from regulatory rules, enabling independent updates to vision models and compliance rule sets without system-wide changes",
    "Deployment: Web-based application supporting desktop, tablet and mobile workflows for field inspection scenarios"
]

add_bullet_points(slide3, bullets_slide3, start_y=Inches(1.8), font_size=Pt(18))

add_footer(slide3, 3)


# ══════════════════════════════════════════════════════════════
# SLIDE 4: FEASIBILITY AND VIABILITY (matches template exactly)
# ══════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(blank_layout)

add_team_name_oval(slide4)
add_sih_logo_top_right(slide4)
add_title(slide4, "FEASIBILITY AND VIABILITY")

bullets_slide4 = [
    "Technology Feasibility: OCR, computer vision and AI vision models (Gemini) are technically mature and applicable to label inspection tasks with documented accuracy in scene-text recognition.",
    "Data Feasibility: Product-label images and structured declaration fields (MRP, manufacturer, net quantity) can be systematically collected for development, testing and continuous improvement.",
    "Architecture Feasibility: AI extraction and regulatory compliance rules are separated into modular components, allowing independent updates to vision models and rule definitions.",
    "Deployment Feasibility: A web-based system (React + Convex) supports desktop, tablet and mobile workflows required for field inspection scenarios.",
    "Scalability: New product categories, declaration types and rule sets can be added through configurable definitions without code changes to the core system.",
    "Challenges and Mitigations: AI misreading text → Confidence scoring + evidence + inspector correction • Blurry images → Image-quality detection + adaptive recapture • Declarations on other sides → Multi-view package analysis • Online information differs → Cross-source comparison • Regulations change → Version-controlled rule engine • False AI findings → Human-in-the-loop verification"
]

add_bullet_points(slide4, bullets_slide4, start_y=Inches(1.8), font_size=Pt(18))

add_footer(slide4, 4)


# ══════════════════════════════════════════════════════════════
# SLIDE 5: IMPACT AND BENEFITS (matches template exactly)
# ══════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(blank_layout)

add_team_name_oval(slide5)
add_sih_logo_top_right(slide5)
add_title(slide5, "IMPACT AND BENEFITS")

bullets_slide5 = [
    "Potential Impact on Target Audience: Faster inspection through automated label-reading and initial validation; consistent checking via systematic rule-based validation; better evidence through image-linked findings; reduced false positives by distinguishing missing from unreadable declarations; improved decision support with recommended next actions for inspectors.",
    "Inspector and Operational Benefits: AI-assisted extraction reduces repetitive manual work; evidence-based findings provide traceable inspection records; risk prioritization focuses inspector attention on high-priority violations; digital reports replace manual documentation; human-in-the-loop ensures inspector retains final authority.",
    "Economic and Strategic Benefits: Reduced manual effort in repetitive label verification; scalable workflow supporting larger inspection volumes; structured digital documentation for audit and compliance records; data-driven monitoring enabling analysis of recurring findings across inspections; future expansion to additional product categories and rule sets.",
    "Social Benefits: Improved consumer protection through more consistent compliance checking; enhanced transparency in inspection processes; support for Legal Metrology enforcement at scale; digital transformation of government inspection workflows."
]

add_bullet_points(slide5, bullets_slide5, start_y=Inches(1.8), font_size=Pt(18))

add_footer(slide5, 5)


# ══════════════════════════════════════════════════════════════
# SLIDE 6: RESEARCH AND REFERENCES (matches template exactly)
# ══════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(blank_layout)

add_team_name_oval(slide6)
add_sih_logo_top_right(slide6)
add_title(slide6, "RESEARCH AND REFERENCES")

bullets_slide6 = [
    "Legal and Regulatory References: Legal Metrology Act, 2009 • Legal Metrology (Packaged Commodities) Rules, 2011 • Department of Consumer Affairs — Legal Metrology Division • Relevant official notifications and amendments under the Legal Metrology framework",
    "Research Areas: Computer Vision and OCR for scene-text recognition in natural environments • Multimodal AI for combined image and text understanding • Explainable AI for confidence scoring and evidence-based decision support • Human-AI Interaction for verified inspection workflows",
    "Technology References: React (UI library) • TypeScript (type-safe development) • Convex (serverless backend and database) • Tailwind CSS (styling) • Gemini AI Vision (image analysis and OCR) • Vite (build tool) • Recharts (analytics visualization)",
    "References: [ADD VERIFIED REFERENCE — Legal Metrology Act, 2009] • [ADD VERIFIED REFERENCE — Legal Metrology (Packaged Commodities) Rules, 2011] • [ADD VERIFIED REFERENCE — Department of Consumer Affairs, Government of India] • [ADD VERIFIED REFERENCE — Relevant court judgments or regulatory orders]"
]

add_bullet_points(slide6, bullets_slide6, start_y=Inches(1.8), font_size=Pt(18))

add_footer(slide6, 6)


# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
output_path = "/home/daytona/codebase/MetrologyAI_SIH_2026_Final.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Slides: {len(prs.slides)}")
print(f"   Format: Widescreen 16:9 (matching SIH template)")
